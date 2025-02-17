import g4f
from pptx import Presentation
import random 
import re
from pptx.util import Pt  # For font size
from pptx.enum.text import PP_ALIGN  # For text alignment
from pptx.dml.color import RGBColor  # For colors
Prompt = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.
Notice
-You do all the presentation text for the user.
-You write the texts no longer than 250 characters!
-You make very short titles!
-You make the presentation easy to understand.
-The presentation has a table of contents.
-The presentation has a summary.
-At least 7 slides.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END"""

def create_ppt_text(Input):
    response = g4f.ChatCompletion.create(
        model="gpt-4",
        provider=g4f.Provider.Yqcloud,
        messages=[
            {"role": "system", "content": (Prompt)},
            {"role": "user", "content": ("The user wants a presentation about " + Input)}
        ],
        stream=True,
    )
    
    ms=""
    for message in response:
        ms+=str(message)
        print(message, end="", flush=True)
    print()
    
    return ms

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation()
    
    with open(text_file, "r", encoding="utf-8") as file:
        lines = file.readlines()
    
    slides_content = []
    title = None
    content = ""

    for line in lines:
        line = line.strip().replace("**", "")  # Remove "**" symbols
        if line.startswith("###"):
            if title:
                slides_content.append((title, content.strip()))
            title = line.replace("###", "").strip()
            content = ""
        else:
            content += line + "\n"
    
    if title:
        slides_content.append((title, content.strip()))
    
    for title, content in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title Formatting (Bold & Larger Font)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        title_text_frame = title_placeholder.text_frame
        title_text_frame.paragraphs[0].font.bold = True
        title_text_frame.paragraphs[0].font.size = Pt(32)  # Larger font size

        # Content Formatting (Bullet Points & Smaller Font)
        content_placeholder = slide.placeholders[1]
        content_text_frame = content_placeholder.text_frame
        content_text_frame.clear()  # Remove default text
        
        for bullet in content.split("\n"):
            if bullet.strip():
                p = content_text_frame.add_paragraph()
                p.text = f"• {bullet.strip()}"  # Add bullet points
                p.space_after = Pt(5)
                p.font.size = Pt(18)  # Smaller font for readability
                p.font.color.rgb = RGBColor(50, 50, 50)  # Dark gray color
                p.alignment = PP_ALIGN.LEFT  # Align text to left

    prs.save(f'POWERPOINT/GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"POWERPOINT/GeneratedPresentations/{ppt_name}.pptx"
    
    return f"{file_path}"

def get_bot_response(msg):
    user_text = msg
    last_char = user_text[-1]
    input_string = user_text
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
    input_string = input_string.replace("\n", "")
    number = 1

    if last_char.isdigit():
        number = int(last_char)
        input_string = user_text[:-2]
        print("Design Number:", number, "selected.")
    else:
        print("No design specified, using default design...")
        
    if number > 7:
        number = 1
        print("Unavailable design, using default design...")
    elif number == 0:
        number = 1
        print("Unavailable design, using default design...")

    with open(f'POWERPOINT/Cache/{input_string}.txt', 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(input_string))

    pptlink = create_ppt(f'POWERPOINT/Cache/{input_string}.txt', number, input_string)
    return str(pptlink)

